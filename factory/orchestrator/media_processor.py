"""
Processador de Midia - Plataforma E
=========================================

Processa videos, audios e documentos para extracao de requisitos.
Utiliza multiplos agentes colaborando para analise completa.
"""

import os
import json
import subprocess
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass, field
import threading
import queue

# Imports da Fabrica
import sys
sys.path.insert(0, str(Path(__file__).parent.parent.parent))


@dataclass
class MediaFile:
    """Representa um arquivo de midia"""
    file_path: str
    file_type: str  # video, audio, document
    file_name: str
    file_size: int
    status: str = "pending"  # pending, processing, completed, error
    extracted_content: Optional[str] = None
    metadata: Dict = field(default_factory=dict)
    processing_started: Optional[datetime] = None
    processing_completed: Optional[datetime] = None
    error_message: Optional[str] = None


@dataclass
class ExtractionResult:
    """Resultado da extracao de conteudo"""
    source_file: str
    file_type: str
    text_content: str
    timestamps: List[Dict] = field(default_factory=list)  # Para videos/audios
    frames: List[str] = field(default_factory=list)  # Caminhos de frames extraidos
    diagrams: List[Dict] = field(default_factory=list)  # Diagramas detectados
    requirements: List[str] = field(default_factory=list)  # Requisitos extraidos
    metadata: Dict = field(default_factory=dict)


class MediaProcessor:
    """
    Processador de Midia

    Coordena multiplos agentes para:
    - Transcrever audios e videos (Whisper)
    - Extrair frames de videos
    - OCR em imagens/slides
    - Extrair texto de documentos
    - Identificar diagramas e fluxos
    """

    def __init__(self, output_dir: str = None):
        self.output_dir = Path(output_dir) if output_dir else Path.cwd() / "media_output"
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.processing_queue: queue.Queue = queue.Queue()
        self.results: Dict[str, ExtractionResult] = {}

        # Status de processamento
        self.is_processing = False
        self._worker_thread: Optional[threading.Thread] = None

        # Callbacks
        self.on_progress = None
        self.on_complete = None
        self.on_error = None

    def scan_directory(self, source_path: str) -> List[MediaFile]:
        """Escaneia diretorio e identifica arquivos de midia"""
        source = Path(source_path)
        if not source.exists():
            return []

        media_files = []

        # Extensoes suportadas
        video_ext = {'.mp4', '.avi', '.mov', '.mkv', '.webm'}
        audio_ext = {'.mp3', '.wav', '.m4a', '.ogg', '.flac'}
        doc_ext = {'.docx', '.doc', '.pdf', '.txt', '.md', '.xlsx', '.pptx'}

        for file_path in source.rglob('*'):
            if file_path.is_file():
                ext = file_path.suffix.lower()

                if ext in video_ext:
                    file_type = "video"
                elif ext in audio_ext:
                    file_type = "audio"
                elif ext in doc_ext:
                    file_type = "document"
                else:
                    continue

                media_files.append(MediaFile(
                    file_path=str(file_path),
                    file_type=file_type,
                    file_name=file_path.name,
                    file_size=file_path.stat().st_size,
                    metadata={"extension": ext}
                ))

        return media_files

    def process_file(self, media_file: MediaFile) -> ExtractionResult:
        """Processa um arquivo de midia"""
        media_file.status = "processing"
        media_file.processing_started = datetime.now()

        try:
            if media_file.file_type == "document":
                result = self._process_document(media_file)
            elif media_file.file_type == "video":
                result = self._process_video(media_file)
            elif media_file.file_type == "audio":
                result = self._process_audio(media_file)
            else:
                raise ValueError(f"Tipo de arquivo nao suportado: {media_file.file_type}")

            media_file.status = "completed"
            media_file.processing_completed = datetime.now()
            media_file.extracted_content = result.text_content[:1000] if result.text_content else None

            return result

        except Exception as e:
            media_file.status = "error"
            media_file.error_message = str(e)
            media_file.processing_completed = datetime.now()

            if self.on_error:
                self.on_error(media_file, e)

            raise

    def _process_document(self, media_file: MediaFile) -> ExtractionResult:
        """Processa documento (DOCX, PDF, TXT, etc)"""
        file_path = Path(media_file.file_path)
        ext = file_path.suffix.lower()
        text_content = ""

        if ext == '.docx':
            text_content = self._extract_docx(file_path)
        elif ext == '.txt' or ext == '.md':
            text_content = file_path.read_text(encoding='utf-8')
        elif ext == '.pdf':
            text_content = self._extract_pdf(file_path)
        elif ext == '.pptx':
            text_content = self._extract_pptx(file_path)

        # Extrai requisitos do texto
        requirements = self._extract_requirements(text_content)

        return ExtractionResult(
            source_file=str(file_path),
            file_type="document",
            text_content=text_content,
            requirements=requirements,
            metadata={"pages": text_content.count('\n\n') + 1}
        )

    def _process_video(self, media_file: MediaFile) -> ExtractionResult:
        """Processa video - extrai audio e frames"""
        file_path = Path(media_file.file_path)
        video_output_dir = self.output_dir / file_path.stem

        video_output_dir.mkdir(parents=True, exist_ok=True)

        # 1. Extrair audio do video
        audio_path = video_output_dir / "audio.wav"
        self._extract_audio_from_video(file_path, audio_path)

        # 2. Transcrever audio
        transcription = ""
        timestamps = []

        if audio_path.exists():
            transcription, timestamps = self._transcribe_audio(audio_path)

        # 3. Extrair frames chave (a cada 30 segundos)
        frames = self._extract_key_frames(file_path, video_output_dir)

        # 4. OCR nos frames para detectar texto/diagramas
        diagrams = []
        for frame_path in frames:
            ocr_text = self._ocr_image(frame_path)
            if ocr_text:
                diagrams.append({
                    "frame": str(frame_path),
                    "text": ocr_text,
                    "type": "slide" if len(ocr_text) > 50 else "diagram"
                })

        # 5. Extrair requisitos
        requirements = self._extract_requirements(transcription)

        return ExtractionResult(
            source_file=str(file_path),
            file_type="video",
            text_content=transcription,
            timestamps=timestamps,
            frames=[str(f) for f in frames],
            diagrams=diagrams,
            requirements=requirements,
            metadata={
                "duration": self._get_video_duration(file_path),
                "frames_extracted": len(frames)
            }
        )

    def _process_audio(self, media_file: MediaFile) -> ExtractionResult:
        """Processa arquivo de audio"""
        file_path = Path(media_file.file_path)

        # Transcrever audio
        transcription, timestamps = self._transcribe_audio(file_path)

        # Extrair requisitos
        requirements = self._extract_requirements(transcription)

        return ExtractionResult(
            source_file=str(file_path),
            file_type="audio",
            text_content=transcription,
            timestamps=timestamps,
            requirements=requirements,
            metadata={
                "duration": self._get_audio_duration(file_path)
            }
        )

    def _extract_docx(self, file_path: Path) -> str:
        """Extrai texto de arquivo DOCX"""
        try:
            from docx import Document
            doc = Document(str(file_path))
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return '\n\n'.join(paragraphs)
        except ImportError:
            return f"[Erro: python-docx nao instalado. Instale com: pip install python-docx]"
        except Exception as e:
            return f"[Erro ao ler DOCX: {e}]"

    def _extract_pdf(self, file_path: Path) -> str:
        """Extrai texto de arquivo PDF"""
        try:
            import PyPDF2
            text = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text.append(page.extract_text())
            return '\n\n'.join(text)
        except ImportError:
            return f"[Erro: PyPDF2 nao instalado. Instale com: pip install PyPDF2]"
        except Exception as e:
            return f"[Erro ao ler PDF: {e}]"

    def _extract_pptx(self, file_path: Path) -> str:
        """Extrai texto de arquivo PowerPoint"""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n\n'.join(text)
        except ImportError:
            return f"[Erro: python-pptx nao instalado. Instale com: pip install python-pptx]"
        except Exception as e:
            return f"[Erro ao ler PPTX: {e}]"

    def _extract_audio_from_video(self, video_path: Path, output_path: Path):
        """Extrai trilha de audio de um video usando ffmpeg"""
        try:
            cmd = [
                'ffmpeg', '-i', str(video_path),
                '-vn', '-acodec', 'pcm_s16le',
                '-ar', '16000', '-ac', '1',
                str(output_path), '-y'
            ]
            subprocess.run(cmd, capture_output=True, check=True)
        except FileNotFoundError:
            print("[Aviso: ffmpeg nao encontrado. Instale ffmpeg para processar videos]")
        except Exception as e:
            print(f"[Erro ao extrair audio: {e}]")

    def _transcribe_audio(self, audio_path: Path) -> Tuple[str, List[Dict]]:
        """Transcreve audio usando Whisper ou similar"""
        try:
            import whisper
            model = whisper.load_model("base")
            result = model.transcribe(str(audio_path))

            text = result.get("text", "")
            segments = result.get("segments", [])

            timestamps = [
                {
                    "start": seg["start"],
                    "end": seg["end"],
                    "text": seg["text"]
                }
                for seg in segments
            ]

            return text, timestamps

        except ImportError:
            # Fallback: retorna placeholder
            return f"[Transcricao pendente - instale whisper: pip install openai-whisper]", []
        except Exception as e:
            return f"[Erro na transcricao: {e}]", []

    def _extract_key_frames(self, video_path: Path, output_dir: Path) -> List[Path]:
        """Extrai frames chave do video"""
        frames = []
        try:
            # Extrai um frame a cada 30 segundos
            cmd = [
                'ffmpeg', '-i', str(video_path),
                '-vf', 'fps=1/30',
                str(output_dir / 'frame_%04d.jpg'), '-y'
            ]
            subprocess.run(cmd, capture_output=True, check=True)

            # Lista frames gerados
            frames = list(output_dir.glob('frame_*.jpg'))

        except FileNotFoundError:
            print("[Aviso: ffmpeg nao encontrado]")
        except Exception as e:
            print(f"[Erro ao extrair frames: {e}]")

        return frames

    def _ocr_image(self, image_path: Path) -> str:
        """Executa OCR em uma imagem"""
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(image_path)
            text = pytesseract.image_to_string(img, lang='por')
            return text.strip()
        except ImportError:
            return ""
        except Exception:
            return ""

    def _get_video_duration(self, video_path: Path) -> float:
        """Obtem duracao do video em segundos"""
        try:
            cmd = [
                'ffprobe', '-v', 'error',
                '-show_entries', 'format=duration',
                '-of', 'default=noprint_wrappers=1:nokey=1',
                str(video_path)
            ]
            result = subprocess.run(cmd, capture_output=True, text=True)
            return float(result.stdout.strip())
        except:
            return 0.0

    def _get_audio_duration(self, audio_path: Path) -> float:
        """Obtem duracao do audio em segundos"""
        return self._get_video_duration(audio_path)

    def _extract_requirements(self, text: str) -> List[str]:
        """Extrai requisitos do texto usando analise basica"""
        requirements = []

        if not text:
            return requirements

        # Palavras-chave que indicam requisitos
        keywords = [
            'precisa', 'deve', 'necessario', 'importante',
            'queremos', 'quero', 'precisamos', 'devemos',
            'obrigatorio', 'essencial', 'funcionalidade',
            'sistema deve', 'usuario deve', 'aplicacao deve'
        ]

        lines = text.split('\n')
        for line in lines:
            line_lower = line.lower()
            for keyword in keywords:
                if keyword in line_lower and len(line) > 20:
                    requirements.append(line.strip())
                    break

        return list(set(requirements))[:50]  # Limita a 50 requisitos

    def process_directory_async(self, source_path: str, callback=None):
        """Processa diretorio em background"""
        def worker():
            self.is_processing = True
            files = self.scan_directory(source_path)

            for i, media_file in enumerate(files):
                try:
                    if self.on_progress:
                        self.on_progress(i + 1, len(files), media_file.file_name)

                    result = self.process_file(media_file)
                    self.results[media_file.file_path] = result

                except Exception as e:
                    print(f"Erro processando {media_file.file_name}: {e}")

            self.is_processing = False

            if self.on_complete:
                self.on_complete(self.results)

            if callback:
                callback(self.results)

        self._worker_thread = threading.Thread(target=worker, daemon=True)
        self._worker_thread.start()

    def get_all_text_content(self) -> str:
        """Retorna todo o conteudo textual extraido"""
        texts = []
        for result in self.results.values():
            if result.text_content:
                texts.append(f"=== {result.source_file} ===\n{result.text_content}")
        return '\n\n'.join(texts)

    def get_all_requirements(self) -> List[str]:
        """Retorna todos os requisitos extraidos"""
        all_reqs = []
        for result in self.results.values():
            all_reqs.extend(result.requirements)
        return list(set(all_reqs))


# Funcao auxiliar para uso rapido
def process_project_media(source_path: str, output_path: str = None) -> Dict:
    """
    Processa todos os arquivos de midia de um projeto

    Args:
        source_path: Caminho dos arquivos fonte
        output_path: Caminho para saida (opcional)

    Returns:
        Dicionario com resultados da extracao
    """
    processor = MediaProcessor(output_path)
    files = processor.scan_directory(source_path)

    results = {}
    for media_file in files:
        try:
            result = processor.process_file(media_file)
            results[media_file.file_path] = {
                "type": result.file_type,
                "text": result.text_content[:500] if result.text_content else "",
                "requirements": result.requirements,
                "metadata": result.metadata
            }
        except Exception as e:
            results[media_file.file_path] = {"error": str(e)}

    return {
        "files_processed": len(files),
        "results": results,
        "all_requirements": processor.get_all_requirements()
    }
