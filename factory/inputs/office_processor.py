"""
Office Processor - Enhanced Microsoft Office Document Processing
Plataforma E v6.5 - Issue #128

Features:
- Word (.docx) document parsing with styles and tables
- Excel (.xlsx) spreadsheet parsing with formulas and charts
- PowerPoint (.pptx) presentation parsing with slides and notes
- Template detection and requirement extraction
- Automatic user story generation with Claude
- Batch processing support
"""

import os
import io
import re
import json
import logging
from typing import Optional, List, Dict, Any, Union, Tuple
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass, field
from enum import Enum

# Microsoft Office document processing
try:
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt
    from docx.enum.style import WD_STYLE_TYPE
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    DocxDocument = None

try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    load_workbook = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    PyPDF2 = None

try:
    from anthropic import Anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False
    Anthropic = None

from factory.config import ANTHROPIC_API_KEY, CLAUDE_MODEL

logger = logging.getLogger(__name__)


class DocumentType(str, Enum):
    """Supported document types"""
    WORD = "docx"
    EXCEL = "xlsx"
    POWERPOINT = "pptx"
    PDF = "pdf"
    TEXT = "txt"
    UNKNOWN = "unknown"


class ExtractionMode(str, Enum):
    """Content extraction modes"""
    FULL = "full"           # Extract all content
    REQUIREMENTS = "requirements"  # Focus on requirements
    SUMMARY = "summary"     # Generate summary only
    STRUCTURED = "structured"  # Preserve structure


@dataclass
class DocumentSection:
    """Represents a section of a document"""
    title: str
    content: str
    level: int = 0
    section_type: str = "text"
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class TableData:
    """Represents table data from a document"""
    headers: List[str]
    rows: List[List[str]]
    title: Optional[str] = None
    source_sheet: Optional[str] = None


@dataclass
class DocumentContent:
    """Extracted content from a document"""
    document_type: DocumentType
    filename: str
    sections: List[DocumentSection] = field(default_factory=list)
    tables: List[TableData] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    raw_text: str = ""
    processing_time: float = 0.0

    def to_dict(self) -> Dict[str, Any]:
        return {
            "document_type": self.document_type.value,
            "filename": self.filename,
            "sections_count": len(self.sections),
            "tables_count": len(self.tables),
            "images_count": len(self.images),
            "metadata": self.metadata,
            "text_length": len(self.raw_text),
            "processing_time": self.processing_time
        }

    def get_full_text(self) -> str:
        """Get combined text from all sections"""
        parts = []
        for section in self.sections:
            if section.title:
                parts.append(f"\n{'#' * (section.level + 1)} {section.title}\n")
            parts.append(section.content)
        return "\n".join(parts)


class OfficeProcessorError(Exception):
    """Custom exception for Office processing errors"""
    pass


class OfficeProcessor:
    """
    Enhanced processor for Microsoft Office documents.

    Features:
    - Word documents with styles, headings, and tables
    - Excel spreadsheets with multiple sheets and formulas
    - PowerPoint presentations with slides and speaker notes
    - Automatic requirement detection and story generation
    - Template-based extraction

    Example:
        processor = OfficeProcessor()
        content = await processor.process_document(file_bytes, "requirements.docx")
        stories = await processor.extract_requirements(content)
    """

    # File size limits
    MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
    MAX_CONTENT_LENGTH = 100000  # Characters for Claude

    # Extension mapping
    EXTENSION_MAP = {
        ".docx": DocumentType.WORD,
        ".doc": DocumentType.WORD,
        ".xlsx": DocumentType.EXCEL,
        ".xls": DocumentType.EXCEL,
        ".pptx": DocumentType.POWERPOINT,
        ".ppt": DocumentType.POWERPOINT,
        ".pdf": DocumentType.PDF,
        ".txt": DocumentType.TEXT
    }

    # Requirement patterns
    REQUIREMENT_PATTERNS = [
        r"(?i)requisito\s*[:\d\.]",
        r"(?i)rf\s*\d+",
        r"(?i)rnf\s*\d+",
        r"(?i)como\s+um?\s+\w+.*quero",
        r"(?i)o\s+sistema\s+deve",
        r"(?i)deve\s+ser\s+possivel",
        r"(?i)o\s+usuario\s+pode",
        r"(?i)funcionalidade\s*:",
        r"(?i)caso\s+de\s+uso",
        r"(?i)historia\s+de\s+usuario"
    ]

    def __init__(
        self,
        claude_model: Optional[str] = None,
        extraction_mode: ExtractionMode = ExtractionMode.REQUIREMENTS
    ):
        """
        Initialize OfficeProcessor.

        Args:
            claude_model: Claude model for story extraction
            extraction_mode: Default extraction mode
        """
        self.claude_model = claude_model or CLAUDE_MODEL
        self.extraction_mode = extraction_mode
        self.claude = None
        self._initialized = False

        # Stats
        self._stats = {
            "documents_processed": 0,
            "stories_extracted": 0,
            "by_type": {t.value: 0 for t in DocumentType}
        }

    def _ensure_initialized(self) -> None:
        """Initialize Claude client."""
        if self._initialized:
            return

        if ANTHROPIC_AVAILABLE and ANTHROPIC_API_KEY:
            self.claude = Anthropic(api_key=ANTHROPIC_API_KEY)
            logger.info("OfficeProcessor Claude initialized")

        self._initialized = True

    def detect_document_type(self, filename: str) -> DocumentType:
        """Detect document type from filename."""
        ext = Path(filename).suffix.lower()
        return self.EXTENSION_MAP.get(ext, DocumentType.UNKNOWN)

    def validate_file(self, file_data: bytes, filename: str) -> DocumentType:
        """Validate file and return document type."""
        if len(file_data) > self.MAX_FILE_SIZE:
            raise OfficeProcessorError(
                f"File too large. Max: {self.MAX_FILE_SIZE / (1024*1024):.0f}MB"
            )

        doc_type = self.detect_document_type(filename)
        if doc_type == DocumentType.UNKNOWN:
            raise OfficeProcessorError(
                f"Unsupported format: {Path(filename).suffix}"
            )

        # Check library availability
        if doc_type == DocumentType.WORD and not DOCX_AVAILABLE:
            raise OfficeProcessorError(
                "python-docx not installed. Run: pip install python-docx"
            )
        if doc_type == DocumentType.EXCEL and not XLSX_AVAILABLE:
            raise OfficeProcessorError(
                "openpyxl not installed. Run: pip install openpyxl"
            )
        if doc_type == DocumentType.POWERPOINT and not PPTX_AVAILABLE:
            raise OfficeProcessorError(
                "python-pptx not installed. Run: pip install python-pptx"
            )
        if doc_type == DocumentType.PDF and not PDF_AVAILABLE:
            raise OfficeProcessorError(
                "PyPDF2 not installed. Run: pip install PyPDF2"
            )

        return doc_type

    def _extract_docx(self, file_data: bytes, filename: str) -> DocumentContent:
        """Extract content from Word document."""
        doc = DocxDocument(io.BytesIO(file_data))
        content = DocumentContent(
            document_type=DocumentType.WORD,
            filename=filename
        )

        # Extract document properties
        core_props = doc.core_properties
        content.metadata = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "created": core_props.created.isoformat() if core_props.created else None,
            "modified": core_props.modified.isoformat() if core_props.modified else None,
            "sections_count": len(doc.sections)
        }

        # Extract paragraphs and headings
        current_section = None
        raw_text_parts = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else "Normal"

            # Check if heading
            if style_name.startswith("Heading"):
                level = int(style_name.replace("Heading ", "").replace("Heading", "1"))
                current_section = DocumentSection(
                    title=text,
                    content="",
                    level=level,
                    section_type="heading"
                )
                content.sections.append(current_section)
            else:
                if current_section:
                    current_section.content += text + "\n"
                else:
                    content.sections.append(DocumentSection(
                        title="",
                        content=text,
                        level=0,
                        section_type="paragraph"
                    ))

            raw_text_parts.append(text)

        # Extract tables
        for table in doc.tables:
            headers = []
            rows = []

            for i, row in enumerate(table.rows):
                row_data = [cell.text.strip() for cell in row.cells]

                if i == 0:
                    headers = row_data
                else:
                    rows.append(row_data)

                raw_text_parts.append(" | ".join(row_data))

            content.tables.append(TableData(
                headers=headers,
                rows=rows
            ))

        content.raw_text = "\n".join(raw_text_parts)
        return content

    def _extract_xlsx(self, file_data: bytes, filename: str) -> DocumentContent:
        """Extract content from Excel spreadsheet."""
        wb = load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)
        content = DocumentContent(
            document_type=DocumentType.EXCEL,
            filename=filename
        )

        content.metadata = {
            "sheets": wb.sheetnames,
            "sheet_count": len(wb.sheetnames)
        }

        raw_text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]

            # Create section for each sheet
            section = DocumentSection(
                title=sheet_name,
                content="",
                level=1,
                section_type="sheet"
            )

            # Extract data
            headers = []
            rows = []
            sheet_text = []

            for row_idx, row in enumerate(sheet.iter_rows()):
                row_values = []
                for cell in row:
                    value = str(cell.value) if cell.value is not None else ""
                    row_values.append(value)

                row_text = " | ".join(row_values)
                sheet_text.append(row_text)

                if row_idx == 0:
                    headers = row_values
                else:
                    rows.append(row_values)

            section.content = "\n".join(sheet_text)
            content.sections.append(section)

            if headers or rows:
                content.tables.append(TableData(
                    headers=headers,
                    rows=rows,
                    source_sheet=sheet_name
                ))

            raw_text_parts.extend(sheet_text)

        wb.close()
        content.raw_text = "\n".join(raw_text_parts)
        return content

    def _extract_pptx(self, file_data: bytes, filename: str) -> DocumentContent:
        """Extract content from PowerPoint presentation."""
        prs = Presentation(io.BytesIO(file_data))
        content = DocumentContent(
            document_type=DocumentType.POWERPOINT,
            filename=filename
        )

        content.metadata = {
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height
        }

        raw_text_parts = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_title = ""
            slide_content = []

            for shape in slide.shapes:
                # Title
                if shape.has_text_frame:
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        if shape.placeholder_format.type == 1:  # Title
                            slide_title = shape.text.strip()
                            continue

                    text = shape.text.strip()
                    if text:
                        slide_content.append(text)

                # Tables
                if shape.has_table:
                    headers = []
                    rows = []
                    for row_idx, row in enumerate(shape.table.rows):
                        row_data = [cell.text.strip() for cell in row.cells]
                        if row_idx == 0:
                            headers = row_data
                        else:
                            rows.append(row_data)
                        slide_content.append(" | ".join(row_data))

                    content.tables.append(TableData(
                        headers=headers,
                        rows=rows,
                        title=f"Slide {slide_idx} Table"
                    ))

            # Speaker notes
            notes = ""
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes = notes_frame.text.strip()

            section = DocumentSection(
                title=slide_title or f"Slide {slide_idx}",
                content="\n".join(slide_content),
                level=1,
                section_type="slide",
                metadata={"slide_number": slide_idx, "notes": notes}
            )
            content.sections.append(section)

            raw_text_parts.append(f"=== Slide {slide_idx}: {slide_title} ===")
            raw_text_parts.extend(slide_content)
            if notes:
                raw_text_parts.append(f"[Notes: {notes}]")

        content.raw_text = "\n".join(raw_text_parts)
        return content

    def _extract_pdf(self, file_data: bytes, filename: str) -> DocumentContent:
        """Extract content from PDF document."""
        reader = PyPDF2.PdfReader(io.BytesIO(file_data))
        content = DocumentContent(
            document_type=DocumentType.PDF,
            filename=filename
        )

        content.metadata = {
            "page_count": len(reader.pages),
            "pdf_metadata": dict(reader.metadata) if reader.metadata else {}
        }

        raw_text_parts = []

        for page_idx, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""

            section = DocumentSection(
                title=f"Page {page_idx}",
                content=page_text.strip(),
                level=1,
                section_type="page"
            )
            content.sections.append(section)

            raw_text_parts.append(f"=== Page {page_idx} ===")
            raw_text_parts.append(page_text)

        content.raw_text = "\n".join(raw_text_parts)
        return content

    def _extract_txt(self, file_data: bytes, filename: str) -> DocumentContent:
        """Extract content from plain text file."""
        # Try different encodings
        text = None
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                text = file_data.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        if text is None:
            text = file_data.decode("utf-8", errors="replace")

        content = DocumentContent(
            document_type=DocumentType.TEXT,
            filename=filename,
            raw_text=text
        )

        # Split by double newlines to create sections
        paragraphs = text.split("\n\n")
        for para in paragraphs:
            para = para.strip()
            if para:
                content.sections.append(DocumentSection(
                    title="",
                    content=para,
                    level=0,
                    section_type="paragraph"
                ))

        content.metadata = {
            "line_count": text.count("\n") + 1,
            "char_count": len(text)
        }

        return content

    async def parse_document(
        self,
        file_data: bytes,
        filename: str
    ) -> DocumentContent:
        """
        Parse a document and extract structured content.

        Args:
            file_data: Raw file bytes
            filename: Document filename

        Returns:
            DocumentContent with extracted content
        """
        import time
        start_time = time.time()

        doc_type = self.validate_file(file_data, filename)

        extractors = {
            DocumentType.WORD: self._extract_docx,
            DocumentType.EXCEL: self._extract_xlsx,
            DocumentType.POWERPOINT: self._extract_pptx,
            DocumentType.PDF: self._extract_pdf,
            DocumentType.TEXT: self._extract_txt
        }

        extractor = extractors.get(doc_type)
        if not extractor:
            raise OfficeProcessorError(f"No extractor for type: {doc_type}")

        try:
            content = extractor(file_data, filename)
            content.processing_time = time.time() - start_time

            self._stats["documents_processed"] += 1
            self._stats["by_type"][doc_type.value] += 1

            logger.info(
                f"Parsed {doc_type.value}: {len(content.sections)} sections, "
                f"{len(content.tables)} tables, {len(content.raw_text)} chars"
            )

            return content

        except Exception as e:
            logger.error(f"Document parsing failed: {e}")
            raise OfficeProcessorError(f"Failed to parse {doc_type.value}: {e}")

    def detect_requirements(self, text: str) -> List[Dict[str, Any]]:
        """
        Detect potential requirements in text using patterns.

        Args:
            text: Text to analyze

        Returns:
            List of detected requirement snippets
        """
        requirements = []

        for pattern in self.REQUIREMENT_PATTERNS:
            matches = re.finditer(pattern, text)
            for match in matches:
                # Get context around match
                start = max(0, match.start() - 50)
                end = min(len(text), match.end() + 200)
                context = text[start:end].strip()

                requirements.append({
                    "pattern": pattern,
                    "match": match.group(),
                    "context": context,
                    "position": match.start()
                })

        return requirements

    async def extract_requirements(
        self,
        content: DocumentContent,
        context: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Extract User Stories from document content using Claude.

        Args:
            content: Parsed document content
            context: Optional project context

        Returns:
            List of extracted User Stories
        """
        self._ensure_initialized()

        if not self.claude:
            raise OfficeProcessorError("Claude not available for extraction")

        # Prepare text for Claude
        text = content.raw_text
        if len(text) > self.MAX_CONTENT_LENGTH:
            text = text[:self.MAX_CONTENT_LENGTH]
            logger.warning("Content truncated for Claude")

        # Include tables in prompt
        tables_text = ""
        for table in content.tables[:5]:  # Limit tables
            table_str = " | ".join(table.headers) + "\n"
            for row in table.rows[:20]:  # Limit rows
                table_str += " | ".join(row) + "\n"
            tables_text += f"\nTable:\n{table_str}"

        prompt = f"""Analise o seguinte documento {content.document_type.value.upper()} e extraia
requisitos de software em formato de User Stories JSON.

Documento: {content.filename}
{f'Contexto: {context}' if context else ''}

Para cada requisito, crie uma User Story com:
- title: Titulo curto e descritivo
- persona: Quem e o usuario
- action: O que ele quer fazer
- benefit: O beneficio esperado
- acceptance_criteria: Lista de criterios de aceite (minimo 3)
- story_points: Estimativa fibonacci (1, 2, 3, 5, 8, 13, 21)
- priority: Prioridade (low, medium, high, urgent)
- category: Categoria (feature, bug, tech_debt, improvement)
- source_location: Onde foi encontrado no documento

Conteudo do documento:
---
{text}
---

{f'Tabelas:{tables_text}' if tables_text else ''}

Responda APENAS com JSON valido:
{{
    "stories": [
        {{
            "title": "...",
            "persona": "...",
            "action": "...",
            "benefit": "...",
            "acceptance_criteria": ["...", "...", "..."],
            "story_points": 5,
            "priority": "medium",
            "category": "feature",
            "source_location": "secao/pagina/slide"
        }}
    ],
    "summary": "resumo do documento",
    "document_type_detected": "tipo identificado",
    "total_requirements_found": 0
}}"""

        try:
            response = self.claude.messages.create(
                model=self.claude_model,
                max_tokens=4096,
                messages=[{"role": "user", "content": prompt}]
            )

            response_text = response.content[0].text.strip()

            # Extract JSON
            if "```json" in response_text:
                json_start = response_text.find("```json") + 7
                json_end = response_text.find("```", json_start)
                response_text = response_text[json_start:json_end].strip()
            elif "```" in response_text:
                json_start = response_text.find("```") + 3
                json_end = response_text.find("```", json_start)
                response_text = response_text[json_start:json_end].strip()

            result = json.loads(response_text)
            stories = result.get("stories", [])

            # Add metadata
            for story in stories:
                story["source"] = "document"
                story["document_type"] = content.document_type.value
                story["filename"] = content.filename
                story["extracted_at"] = datetime.utcnow().isoformat()

            self._stats["stories_extracted"] += len(stories)

            logger.info(f"Extracted {len(stories)} stories from {content.filename}")
            return stories

        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse Claude response: {e}")
            return [{
                "title": f"Documento {content.filename} a revisar",
                "persona": "analista",
                "action": "revisar requisitos do documento",
                "benefit": "processar corretamente os requisitos",
                "acceptance_criteria": ["Revisar documento original"],
                "story_points": 1,
                "priority": "medium",
                "source": "document",
                "document_type": content.document_type.value,
                "needs_review": True,
                "extracted_at": datetime.utcnow().isoformat()
            }]
        except Exception as e:
            logger.error(f"Story extraction failed: {e}")
            raise OfficeProcessorError(f"Failed to extract stories: {e}")

    async def process_document(
        self,
        file_data: bytes,
        filename: str,
        context: Optional[str] = None,
        extract_stories: bool = True
    ) -> Dict[str, Any]:
        """
        Complete pipeline: parse document and extract stories.

        Args:
            file_data: Raw file bytes
            filename: Document filename
            context: Optional project context
            extract_stories: Whether to extract stories

        Returns:
            Dictionary with parsed content and stories
        """
        import time
        start_time = time.time()

        # Parse document
        content = await self.parse_document(file_data, filename)

        # Detect requirements
        detected_reqs = self.detect_requirements(content.raw_text)

        # Extract stories if requested
        stories = []
        if extract_stories:
            stories = await self.extract_requirements(content, context)

        processing_time = time.time() - start_time

        return {
            "content": content.to_dict(),
            "text": content.raw_text,
            "sections": [
                {
                    "title": s.title,
                    "content": s.content[:500],  # Truncate for response
                    "type": s.section_type
                }
                for s in content.sections
            ],
            "tables": [
                {
                    "headers": t.headers,
                    "row_count": len(t.rows),
                    "source": t.source_sheet
                }
                for t in content.tables
            ],
            "detected_requirements": detected_reqs[:10],  # Limit
            "stories": stories,
            "stories_count": len(stories),
            "processing_time": round(processing_time, 2),
            "source": "document",
            "processed_at": datetime.utcnow().isoformat()
        }

    async def batch_process(
        self,
        files: List[Tuple[bytes, str]],
        context: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Process multiple documents in batch.

        Args:
            files: List of (file_data, filename) tuples
            context: Optional shared context

        Returns:
            List of processing results
        """
        results = []
        for file_data, filename in files:
            try:
                result = await self.process_document(file_data, filename, context)
                results.append(result)
            except Exception as e:
                results.append({
                    "filename": filename,
                    "error": str(e),
                    "status": "failed"
                })

        return results

    def get_supported_formats(self) -> Dict[str, bool]:
        """Return supported formats and availability."""
        return {
            "docx": DOCX_AVAILABLE,
            "xlsx": XLSX_AVAILABLE,
            "pptx": PPTX_AVAILABLE,
            "pdf": PDF_AVAILABLE,
            "txt": True
        }

    def get_status(self) -> Dict[str, Any]:
        """Get current status."""
        return {
            "initialized": self._initialized,
            "claude_available": bool(self.claude),
            "supported_formats": self.get_supported_formats(),
            "extraction_mode": self.extraction_mode.value,
            "stats": self._stats,
            "libraries": {
                "python-docx": DOCX_AVAILABLE,
                "openpyxl": XLSX_AVAILABLE,
                "python-pptx": PPTX_AVAILABLE,
                "PyPDF2": PDF_AVAILABLE
            }
        }
