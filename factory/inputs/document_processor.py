"""
Document Processor - Office Documents to Requirements
Fabrica de Agentes v6.5 - Issue #128

Processes DOCX, XLSX, PPTX, and PDF files to extract requirements.
Uses specialized libraries for each format and Claude for intelligent extraction.
"""

import os
import io
import json
import logging
from typing import Optional, List, Dict, Any, Union
from datetime import datetime
from pathlib import Path
from enum import Enum

# Document processing libraries
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    DocxDocument = None

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    load_workbook = None

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    PyPDF2 = None

try:
    from anthropic import Anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False
    Anthropic = None

from factory.config import ANTHROPIC_API_KEY, CLAUDE_MODEL

logger = logging.getLogger(__name__)


class DocumentType(str, Enum):
    """Supported document types"""
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    PDF = "pdf"
    TXT = "txt"
    UNKNOWN = "unknown"


class DocumentProcessorError(Exception):
    """Custom exception for document processing errors"""
    pass


class DocumentProcessor:
    """
    Document processor for extracting requirements from Office documents.

    Supports:
    - DOCX: Microsoft Word documents
    - XLSX: Microsoft Excel spreadsheets
    - PPTX: Microsoft PowerPoint presentations
    - PDF: Adobe PDF documents
    - TXT: Plain text files

    Example:
        processor = DocumentProcessor()
        result = await processor.process_document(file_bytes, "requirements.docx")
    """

    # Maximum file size in bytes (50MB)
    MAX_FILE_SIZE = 50 * 1024 * 1024

    # Maximum text content to send to Claude (chars)
    MAX_CONTENT_LENGTH = 50000

    # Supported file extensions
    SUPPORTED_EXTENSIONS = {
        ".docx": DocumentType.DOCX,
        ".doc": DocumentType.DOCX,  # Will try to process as DOCX
        ".xlsx": DocumentType.XLSX,
        ".xls": DocumentType.XLSX,
        ".pptx": DocumentType.PPTX,
        ".ppt": DocumentType.PPTX,
        ".pdf": DocumentType.PDF,
        ".txt": DocumentType.TXT,
    }

    def __init__(self, claude_model: Optional[str] = None):
        """
        Initialize DocumentProcessor.

        Args:
            claude_model: Claude model to use for extraction (optional)
        """
        self.claude_model = claude_model or CLAUDE_MODEL
        self.claude = None
        self._initialized = False

    def _ensure_initialized(self) -> None:
        """Lazily initialize Claude client."""
        if self._initialized:
            return

        if not ANTHROPIC_AVAILABLE:
            raise DocumentProcessorError(
                "Anthropic SDK is not installed. Install with: pip install anthropic"
            )

        if not ANTHROPIC_API_KEY:
            raise DocumentProcessorError(
                "ANTHROPIC_API_KEY is not configured. Set it in environment variables."
            )

        try:
            self.claude = Anthropic(api_key=ANTHROPIC_API_KEY)
        except Exception as e:
            raise DocumentProcessorError(f"Failed to initialize Claude client: {e}")

        self._initialized = True
        logger.info("DocumentProcessor initialized successfully")

    def detect_document_type(self, filename: str) -> DocumentType:
        """
        Detect document type from filename extension.

        Args:
            filename: Name of the file

        Returns:
            DocumentType enum value
        """
        ext = Path(filename).suffix.lower()
        return self.SUPPORTED_EXTENSIONS.get(ext, DocumentType.UNKNOWN)

    def validate_file(
        self,
        file_data: bytes,
        filename: str
    ) -> DocumentType:
        """
        Validate file before processing.

        Args:
            file_data: Raw file bytes
            filename: Name of the file

        Returns:
            DocumentType if valid

        Raises:
            DocumentProcessorError: If validation fails
        """
        # Check file size
        if len(file_data) > self.MAX_FILE_SIZE:
            raise DocumentProcessorError(
                f"File too large. Maximum size is {self.MAX_FILE_SIZE / (1024*1024):.1f}MB"
            )

        # Check file type
        doc_type = self.detect_document_type(filename)
        if doc_type == DocumentType.UNKNOWN:
            ext = Path(filename).suffix.lower()
            supported = ", ".join(self.SUPPORTED_EXTENSIONS.keys())
            raise DocumentProcessorError(
                f"Unsupported file format: {ext}. Supported formats: {supported}"
            )

        # Check if required library is available
        if doc_type == DocumentType.DOCX and not DOCX_AVAILABLE:
            raise DocumentProcessorError(
                "python-docx is not installed. Install with: pip install python-docx"
            )
        if doc_type == DocumentType.XLSX and not XLSX_AVAILABLE:
            raise DocumentProcessorError(
                "openpyxl is not installed. Install with: pip install openpyxl"
            )
        if doc_type == DocumentType.PPTX and not PPTX_AVAILABLE:
            raise DocumentProcessorError(
                "python-pptx is not installed. Install with: pip install python-pptx"
            )
        if doc_type == DocumentType.PDF and not PDF_AVAILABLE:
            raise DocumentProcessorError(
                "PyPDF2 is not installed. Install with: pip install PyPDF2"
            )

        return doc_type

    def extract_text_docx(self, file_data: bytes) -> Dict[str, Any]:
        """
        Extract text from DOCX file.

        Args:
            file_data: Raw DOCX bytes

        Returns:
            Dictionary with extracted content and metadata
        """
        try:
            doc = DocxDocument(io.BytesIO(file_data))

            # Extract paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append({
                        "text": text,
                        "style": para.style.name if para.style else "Normal"
                    })

            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables.append(table_data)

            # Combine text
            full_text = "\n".join(p["text"] for p in paragraphs)

            # Add table content as text
            for table in tables:
                for row in table:
                    full_text += "\n" + " | ".join(row)

            return {
                "text": full_text,
                "paragraphs": len(paragraphs),
                "tables": len(tables),
                "document_properties": {
                    "sections": len(doc.sections),
                }
            }

        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {e}")
            raise DocumentProcessorError(f"Failed to read DOCX file: {e}")

    def extract_text_xlsx(self, file_data: bytes) -> Dict[str, Any]:
        """
        Extract text from XLSX file.

        Args:
            file_data: Raw XLSX bytes

        Returns:
            Dictionary with extracted content and metadata
        """
        try:
            wb = load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)

            sheets_content = []
            full_text = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"=== Sheet: {sheet_name} ==="]

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    if row_values:
                        row_text = " | ".join(row_values)
                        sheet_text.append(row_text)

                sheets_content.append({
                    "name": sheet_name,
                    "rows": len(sheet_text) - 1
                })
                full_text.extend(sheet_text)

            wb.close()

            return {
                "text": "\n".join(full_text),
                "sheets": sheets_content,
                "total_sheets": len(wb.sheetnames)
            }

        except Exception as e:
            logger.error(f"Failed to extract text from XLSX: {e}")
            raise DocumentProcessorError(f"Failed to read XLSX file: {e}")

    def extract_text_pptx(self, file_data: bytes) -> Dict[str, Any]:
        """
        Extract text from PPTX file.

        Args:
            file_data: Raw PPTX bytes

        Returns:
            Dictionary with extracted content and metadata
        """
        try:
            prs = Presentation(io.BytesIO(file_data))

            slides_content = []
            full_text = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== Slide {slide_num} ==="]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())

                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text:
                                slide_text.append(row_text)

                slides_content.append({
                    "slide_number": slide_num,
                    "text_blocks": len(slide_text) - 1
                })
                full_text.extend(slide_text)

            return {
                "text": "\n".join(full_text),
                "slides": slides_content,
                "total_slides": len(prs.slides)
            }

        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {e}")
            raise DocumentProcessorError(f"Failed to read PPTX file: {e}")

    def extract_text_pdf(self, file_data: bytes) -> Dict[str, Any]:
        """
        Extract text from PDF file.

        Args:
            file_data: Raw PDF bytes

        Returns:
            Dictionary with extracted content and metadata
        """
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(file_data))

            pages_content = []
            full_text = []

            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    full_text.append(f"=== Page {page_num} ===")
                    full_text.append(page_text.strip())
                    pages_content.append({
                        "page_number": page_num,
                        "char_count": len(page_text)
                    })

            return {
                "text": "\n".join(full_text),
                "pages": pages_content,
                "total_pages": len(reader.pages),
                "metadata": dict(reader.metadata) if reader.metadata else {}
            }

        except Exception as e:
            logger.error(f"Failed to extract text from PDF: {e}")
            raise DocumentProcessorError(f"Failed to read PDF file: {e}")

    def extract_text_txt(self, file_data: bytes) -> Dict[str, Any]:
        """
        Extract text from plain text file.

        Args:
            file_data: Raw text bytes

        Returns:
            Dictionary with extracted content and metadata
        """
        try:
            # Try different encodings
            for encoding in ["utf-8", "latin-1", "cp1252"]:
                try:
                    text = file_data.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                text = file_data.decode("utf-8", errors="replace")

            lines = text.split("\n")
            non_empty_lines = [l for l in lines if l.strip()]

            return {
                "text": text,
                "total_lines": len(lines),
                "non_empty_lines": len(non_empty_lines),
                "char_count": len(text)
            }

        except Exception as e:
            logger.error(f"Failed to extract text from TXT: {e}")
            raise DocumentProcessorError(f"Failed to read text file: {e}")

    def extract_text(
        self,
        file_data: bytes,
        doc_type: DocumentType
    ) -> Dict[str, Any]:
        """
        Extract text from document based on type.

        Args:
            file_data: Raw file bytes
            doc_type: Document type

        Returns:
            Dictionary with extracted content and metadata
        """
        extractors = {
            DocumentType.DOCX: self.extract_text_docx,
            DocumentType.XLSX: self.extract_text_xlsx,
            DocumentType.PPTX: self.extract_text_pptx,
            DocumentType.PDF: self.extract_text_pdf,
            DocumentType.TXT: self.extract_text_txt,
        }

        extractor = extractors.get(doc_type)
        if not extractor:
            raise DocumentProcessorError(f"No extractor for type: {doc_type}")

        return extractor(file_data)

    async def extract_requirements(
        self,
        content: str,
        doc_type: str,
        context: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Extract User Stories from document text using Claude.

        Args:
            content: Extracted text content
            doc_type: Type of document (for context)
            context: Optional project context

        Returns:
            List of extracted User Stories as dictionaries
        """
        self._ensure_initialized()

        if not content or not content.strip():
            logger.warning("Empty content provided")
            return []

        # Truncate if too long
        if len(content) > self.MAX_CONTENT_LENGTH:
            content = content[:self.MAX_CONTENT_LENGTH]
            logger.warning(f"Content truncated to {self.MAX_CONTENT_LENGTH} chars")

        prompt = f"""Analise o seguinte documento ({doc_type.upper()}) e extraia requisitos
de software em formato de User Stories JSON.

Para cada requisito identificado, crie uma User Story com:
- title: Título curto e descritivo
- persona: Quem é o usuário
- action: O que ele quer fazer
- benefit: O benefício esperado
- acceptance_criteria: Lista de critérios de aceite (pelo menos 3)
- story_points: Estimativa de complexidade (1, 2, 3, 5, 8, 13 ou 21)
- priority: Prioridade (low, medium, high, urgent)
- category: Categoria (feature, bug, tech_debt, improvement)

{f'Contexto do projeto: {context}' if context else ''}

Documento:
---
{content}
---

Responda APENAS com um JSON válido no formato:
{{
    "stories": [
        {{
            "title": "...",
            "persona": "...",
            "action": "...",
            "benefit": "...",
            "acceptance_criteria": ["...", "...", "..."],
            "story_points": 5,
            "priority": "medium",
            "category": "feature",
            "source_section": "seção relevante do documento"
        }}
    ],
    "summary": "resumo geral dos requisitos identificados",
    "document_type_detected": "tipo de documento detectado",
    "suggestions": ["sugestão 1", "sugestão 2"]
}}

Se não conseguir identificar requisitos claros, retorne um array vazio de stories."""

        try:
            logger.info(f"Extracting requirements from {doc_type} document")

            response = self.claude.messages.create(
                model=self.claude_model,
                max_tokens=4096,
                messages=[
                    {
                        "role": "user",
                        "content": prompt
                    }
                ]
            )

            # Parse response
            response_text = response.content[0].text.strip()

            # Handle JSON wrapped in markdown
            if "```json" in response_text:
                json_start = response_text.find("```json") + 7
                json_end = response_text.find("```", json_start)
                response_text = response_text[json_start:json_end].strip()
            elif "```" in response_text:
                json_start = response_text.find("```") + 3
                json_end = response_text.find("```", json_start)
                response_text = response_text[json_start:json_end].strip()

            result = json.loads(response_text)
            stories = result.get("stories", [])

            # Add metadata to each story
            for story in stories:
                story["source"] = "document"
                story["document_type"] = doc_type
                story["extracted_at"] = datetime.utcnow().isoformat()

            logger.info(f"Extracted {len(stories)} stories from document")
            return stories

        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse Claude response as JSON: {e}")
            return [{
                "title": "Documento a revisar",
                "persona": "analista",
                "action": "revisar requisitos do documento",
                "benefit": "processar corretamente os requisitos",
                "acceptance_criteria": ["Revisar documento original"],
                "story_points": 1,
                "priority": "medium",
                "source": "document",
                "document_type": doc_type,
                "needs_review": True,
                "extracted_at": datetime.utcnow().isoformat()
            }]

        except Exception as e:
            logger.error(f"Requirements extraction failed: {e}")
            raise DocumentProcessorError(f"Failed to extract requirements: {e}")

    async def process_document(
        self,
        file_data: bytes,
        filename: str,
        context: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Complete pipeline: document -> text extraction -> user stories.

        Args:
            file_data: Raw file bytes
            filename: Name of the file
            context: Optional project context

        Returns:
            Dictionary with:
            {
                "stories": List of extracted User Stories,
                "text": Extracted text content,
                "document_type": Type of document,
                "metadata": Document metadata,
                "processing_time": Processing time in seconds
            }
        """
        import time
        start_time = time.time()

        try:
            # Validate and detect type
            doc_type = self.validate_file(file_data, filename)

            # Extract text
            extraction = self.extract_text(file_data, doc_type)

            # Extract requirements
            stories = await self.extract_requirements(
                extraction["text"],
                doc_type.value,
                context=context
            )

            processing_time = time.time() - start_time

            # Build metadata without the full text
            metadata = {k: v for k, v in extraction.items() if k != "text"}

            return {
                "stories": stories,
                "text": extraction["text"],
                "document_type": doc_type.value,
                "filename": filename,
                "file_size": len(file_data),
                "metadata": metadata,
                "processing_time": round(processing_time, 2),
                "source": "document",
                "processed_at": datetime.utcnow().isoformat()
            }

        except DocumentProcessorError:
            raise
        except Exception as e:
            logger.error(f"Document processing pipeline failed: {e}")
            raise DocumentProcessorError(f"Document processing failed: {e}")

    def get_supported_formats(self) -> Dict[str, bool]:
        """Return supported formats and their availability."""
        return {
            "docx": DOCX_AVAILABLE,
            "xlsx": XLSX_AVAILABLE,
            "pptx": PPTX_AVAILABLE,
            "pdf": PDF_AVAILABLE,
            "txt": True
        }

    def get_status(self) -> Dict[str, Any]:
        """Get current status of the DocumentProcessor."""
        return {
            "initialized": self._initialized,
            "claude_model": self.claude_model,
            "anthropic_available": ANTHROPIC_AVAILABLE,
            "max_file_size_mb": self.MAX_FILE_SIZE / (1024 * 1024),
            "supported_formats": self.get_supported_formats(),
            "libraries": {
                "python-docx": DOCX_AVAILABLE,
                "openpyxl": XLSX_AVAILABLE,
                "python-pptx": PPTX_AVAILABLE,
                "PyPDF2": PDF_AVAILABLE
            }
        }
